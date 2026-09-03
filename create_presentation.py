from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import urllib.request
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(25, 55, 109)
LIGHT_BLUE = RGBColor(68, 114, 196)
ACCENT_COLOR = RGBColor(237, 125, 49)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(68, 68, 68)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    # Add title text
    title_frame = title_shape.text_frame
    title_frame.vertical_anchor = 1
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_before = Pt(10)
    p.space_after = Pt(10)
    p.alignment = PP_ALIGN.LEFT
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.line_spacing = 1.3

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    title_frame = title_shape.text_frame
    title_frame.vertical_anchor = 1
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)

# Slide 1: Title Slide
add_title_slide(prs, "AI Development", "Building Intelligent Systems in 2026")

# Slide 2: What is AI?
add_content_slide(prs, "What is Artificial Intelligence?", [
    "🤖 AI is the simulation of human intelligence by machines",
    "🧠 Capabilities include learning, reasoning, and problem-solving",
    "📊 Powered by algorithms and vast amounts of data",
    "⚡ Rapidly evolving field with transformative potential",
    "🎯 Applications across healthcare, finance, education, and more"
])

# Slide 3: AI Development Timeline
add_content_slide(prs, "AI Development Timeline", [
    "1950s-1970s: Early AI Research & Logic-based Systems",
    "1980s-1990s: Expert Systems & Machine Learning Emergence",
    "2000s: Data Explosion & Computational Power Growth",
    "2010s: Deep Learning Revolution & Neural Networks",
    "2020s: Large Language Models & Generative AI Boom"
])

# Slide 4: Types of AI
add_two_column_slide(prs, "Types of AI", 
    [
        "Narrow AI (Weak AI):",
        "• Specialized in specific tasks",
        "• All current AI systems",
        "• Chess engines, chatbots",
        "",
        "General AI (Strong AI):",
        "• Theoretically matches human intelligence",
        "• Multi-domain capabilities"
    ],
    [
        "Cognitive Levels:",
        "• Reactive Machines",
        "• Limited Memory",
        "• Theory of Mind",
        "• Self-Aware AI",
        "",
        "Current Focus:",
        "• Improving narrow AI capabilities"
    ]
)

# Slide 5: Machine Learning Fundamentals
add_content_slide(prs, "Machine Learning Fundamentals", [
    "📈 Supervised Learning: Learning from labeled data (classification, regression)",
    "🎲 Unsupervised Learning: Finding patterns in unlabeled data (clustering)",
    "🎮 Reinforcement Learning: Learning through rewards and penalties",
    "🔄 Deep Learning: Using neural networks with multiple layers",
    "🚀 Transfer Learning: Applying knowledge from one task to another"
])

# Slide 6: AI Development Process
add_content_slide(prs, "AI Development Lifecycle", [
    "1️⃣ Problem Definition: Identify the business challenge",
    "2️⃣ Data Collection: Gather relevant training data",
    "3️⃣ Data Preprocessing: Clean, normalize, and prepare data",
    "4️⃣ Model Selection: Choose appropriate algorithms",
    "5️⃣ Training & Validation: Optimize model performance",
    "6️⃣ Testing & Evaluation: Assess accuracy and reliability",
    "7️⃣ Deployment: Push model to production environment"
])

# Slide 7: Key Technologies in AI
add_two_column_slide(prs, "Key AI Technologies",
    [
        "Neural Networks:",
        "• Inspired by brain structure",
        "• Multiple layers process data",
        "• Foundation of deep learning",
        "",
        "Natural Language Processing:",
        "• Understanding human language",
        "• Text analysis and generation"
    ],
    [
        "Computer Vision:",
        "• Image recognition & analysis",
        "• Object detection",
        "• Facial recognition",
        "",
        "Generative Models:",
        "• Create new content",
        "• GANs and Transformers"
    ]
)

# Slide 8: Popular AI Frameworks & Tools
add_content_slide(prs, "AI Development Frameworks & Tools", [
    "🐍 Python: Leading language for AI/ML development",
    "🧠 TensorFlow & PyTorch: Deep learning frameworks",
    "📚 Scikit-learn: Machine learning library",
    "🤗 Hugging Face: NLP models and transformers",
    "⚙️ Keras: User-friendly neural network API",
    "☁️ Cloud Platforms: AWS SageMaker, Google Cloud AI, Azure ML"
])

# Slide 9: Real-World AI Applications
add_two_column_slide(prs, "Real-World AI Applications",
    [
        "Healthcare:",
        "• Disease diagnosis",
        "• Drug discovery",
        "• Patient monitoring",
        "",
        "Finance:",
        "• Fraud detection",
        "• Algorithmic trading",
        "• Risk assessment"
    ],
    [
        "Technology:",
        "• Virtual assistants",
        "• Recommendation systems",
        "• Autonomous vehicles",
        "",
        "Business:",
        "• Process automation",
        "• Predictive analytics",
        "• Customer service"
    ]
)

# Slide 10: Generative AI Revolution
add_content_slide(prs, "Generative AI Revolution", [
    "💬 ChatGPT & Large Language Models: Conversational AI",
    "🎨 DALL-E & Stable Diffusion: Image generation",
    "🎵 Music Generation: Composing original pieces",
    "📝 Code Generation: AI-assisted programming",
    "🚀 Impact: Democratizing AI capabilities for everyone",
    "⚠️ Challenges: Bias, misinformation, ethical concerns"
])

# Slide 11: Challenges in AI Development
add_content_slide(prs, "Challenges in AI Development", [
    "📊 Data Quality: Poor data leads to poor models",
    "⚖️ Bias & Fairness: Models may perpetuate discrimination",
    "🔐 Security: Adversarial attacks and data privacy concerns",
    "💰 Computational Cost: Training large models is expensive",
    "🧪 Interpretability: Understanding AI decision-making",
    "⚖️ Ethical Concerns: Responsibility and accountability"
])

# Slide 12: Best Practices in AI Development
add_content_slide(prs, "Best Practices in AI Development", [
    "✅ Start with clear objectives and success metrics",
    "✅ Invest time in quality data collection and preprocessing",
    "✅ Use version control for models and data",
    "✅ Implement rigorous testing and validation",
    "✅ Monitor model performance in production",
    "✅ Consider ethical implications from the start",
    "✅ Document assumptions and limitations"
])

# Slide 13: Future of AI
add_content_slide(prs, "Future of AI Development", [
    "🎯 Artificial General Intelligence (AGI): On the horizon but challenges remain",
    "🔬 Quantum Computing: Exponential speedup for AI algorithms",
    "🌍 AI for Social Good: Climate, education, healthcare",
    "🤝 Human-AI Collaboration: Augmenting human capabilities",
    "🛡️ Responsible AI: Ethics, safety, and alignment",
    "📱 Edge AI: Running models on local devices"
])

# Slide 14: Getting Started with AI
add_content_slide(prs, "Getting Started with AI Development", [
    "📚 Learn Python and fundamentals of statistics & math",
    "🎓 Take online courses (Coursera, Fast.ai, Andrew Ng's courses)",
    "💻 Practice with Kaggle competitions and projects",
    "📖 Read research papers and stay updated",
    "👥 Join AI communities and collaborate",
    "🚀 Build your own projects from idea to deployment"
])

# Slide 15: Conclusion
add_title_slide(prs, "Thank You!", "The Future of AI is Yours to Create")

# Save presentation
output_path = "AI_Development_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully: {output_path}")
